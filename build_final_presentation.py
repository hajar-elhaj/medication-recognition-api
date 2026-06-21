# -*- coding: utf-8 -*-
"""Builds the FINAL project presentation (standalone) for Project #2.
Dark + cyan theme, English, 18 slides (title + outline + 15 content + thanks).
Run: python build_final_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Palette ─────────────────────────────────────────────────────────────────
BG     = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
CARD   = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
CARD2  = RGBColor(0x24, 0x31, 0x47)
CYAN   = RGBColor(0x22, 0xD3, 0xEE)
CYAN_D = RGBColor(0x0E, 0x74, 0x90)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x94, 0xA3, 0xB8)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
RED    = RGBColor(0xF8, 0x71, 0x71)

FONT   = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


# ── Helpers ───────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _set_font(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def text(s, left, top, width, height, runs, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=None, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (txt, dict) tuples."""
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for txt, opt in para:
            r = p.add_run()
            r.text = txt
            _set_font(r, opt.get("size", 16), opt.get("color", WHITE),
                      opt.get("bold", False), opt.get("italic", False),
                      opt.get("font", FONT))
        if "space_after" in (para[0][1] if para and isinstance(para[0], tuple) else {}):
            p.space_after = Pt(para[0][1]["space_after"])
    return tb


def card(s, left, top, width, height, fill=CARD, radius=0.08):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def circle(s, left, top, diam, fill=CYAN):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diam, diam)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def numbered(s, left, top, num, diam=Inches(0.55), fill=CYAN, color=BG):
    circle(s, left, top, diam, fill)
    text(s, left, top, diam, diam, [[(str(num), {"size": 22, "color": color, "bold": True})]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def title(s, txt, sub=None):
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
         [[(txt, {"size": 32, "color": WHITE, "bold": True})]])
    if sub:
        text(s, Inches(0.72), Inches(1.18), Inches(12), Inches(0.5),
             [[(sub, {"size": 15, "color": CYAN})]])


def kicker(s, txt):
    """Small cyan label top-left for section identity."""
    text(s, Inches(0.72), Inches(0.32), Inches(8), Inches(0.3),
         [[(txt.upper(), {"size": 12, "color": CYAN, "bold": True})]])


def page_num(s, n):
    """Small slide number, bottom-right corner."""
    text(s, Inches(12.5), Inches(7.02), Inches(0.7), Inches(0.35),
         [[(str(n), {"size": 11, "color": MUTED})]], align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
# decorative concentric circles motif (top-right)
circle(s, Inches(10.6), Inches(-1.4), Inches(4.2), CARD)
circle(s, Inches(11.5), Inches(-0.7), Inches(2.4), CARD2)
circle(s, Inches(12.0), Inches(-0.2), Inches(1.2), CYAN_D)

text(s, Inches(0.8), Inches(1.35), Inches(11), Inches(0.4),
     [[("MOHAMMED V UNIVERSITY  •  FACULTY OF SCIENCES, RABAT  •  MASTER IT",
        {"size": 13, "color": CYAN, "bold": True})]])
text(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.8),
     [[("Medication Box", {"size": 50, "color": WHITE, "bold": True})],
      [("Recognition API", {"size": 50, "color": CYAN, "bold": True})]],
     line_spacing=1.0)
text(s, Inches(0.82), Inches(4.15), Inches(11), Inches(0.5),
     [[("Medicine recognition through Computer Vision & multilingual OCR",
        {"size": 17, "color": MUTED})]])

# names band
text(s, Inches(0.82), Inches(5.5), Inches(11), Inches(1.2),
     [[("Authors: ", {"size": 15, "color": MUTED}),
       ("Hajar El Haj", {"size": 15, "color": WHITE, "bold": True}),
       ("   &   ", {"size": 15, "color": MUTED}),
       ("Ikram Belmalhouz", {"size": 15, "color": WHITE, "bold": True})],
      [("Supervisor: ", {"size": 15, "color": MUTED}),
       ("Pr. Abdelhak Mahmoudi", {"size": 15, "color": CYAN, "bold": True}),
       ("   •   Co-supervisor: ", {"size": 15, "color": MUTED}),
       ("Saad Frihi", {"size": 15, "color": CYAN, "bold": True})]],
     line_spacing=1.4)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OUTLINE / PLAN
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Outline")
title(s, "Outline")

plan = [
    ("Introduction & Objectives", "Problem statement and project goals"),
    ("State of the Art", "Comparison of OCR tools and recognition methods"),
    ("Architecture & Pipeline", "From the image to the bilingual result"),
    ("Key Components", "OCR, adaptive OpenCV, reference database"),
    ("API & Web Interface", "REST service and user interface"),
    ("Evaluation & Results", "Testing, metrics and limitations"),
    ("Conclusion & Perspectives", "Achievements and future work"),
]
y0 = Inches(1.65)
for i, (h, d) in enumerate(plan):
    cy = y0 + Emu(int(Inches(0.76) * i))
    card(s, Inches(0.9), cy, Inches(11.5), Inches(0.64))
    numbered(s, Inches(1.12), cy + Inches(0.075), i + 1, diam=Inches(0.49))
    text(s, Inches(1.95), cy, Inches(5.2), Inches(0.64),
         [[(h, {"size": 16, "color": WHITE, "bold": True})]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(6.6), cy, Inches(5.6), Inches(0.64),
         [[(d, {"size": 13, "color": MUTED})]], anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXT & PROBLEM
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "01 — Introduction")
title(s, "Context & Problem")

text(s, Inches(0.72), Inches(1.5), Inches(6.0), Inches(2.2),
     [[("Identifying a medicine from a photo of its box is highly useful in "
        "digital-health applications: patient assistance, pharmacies, stock "
        "control.", {"size": 16, "color": WHITE})],
      [("", {"size": 6})],
      [("But it is a hard problem under real-world conditions.",
        {"size": 16, "color": CYAN, "bold": True})]],
     line_spacing=1.25)

challenges = [
    ("Visual similarity", "Boxes that look alike (same color, shape)"),
    ("Image quality", "Blur, poor lighting, skewed photos"),
    ("Multilingual text", "Packaging in both Arabic AND French"),
    ("Partial visibility", "Box cropped, folded or flipped"),
]
x = Inches(7.0)
y = Inches(1.5)
for i, (h, d) in enumerate(challenges):
    cy = y + Emu(int(Inches(1.32) * i))
    card(s, x, cy, Inches(5.6), Inches(1.12))
    circle(s, x + Inches(0.28), cy + Inches(0.33), Inches(0.46), CYAN)
    text(s, x + Inches(0.28), cy + Inches(0.33), Inches(0.46), Inches(0.46),
         [[("!", {"size": 20, "color": BG, "bold": True})]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.95), cy + Inches(0.16), Inches(4.5), Inches(0.9),
         [[(h, {"size": 15, "color": WHITE, "bold": True})],
          [(d, {"size": 12.5, "color": MUTED})]], line_spacing=1.1)
page_num(s, 3)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OBJECTIVES
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "02 — Objectives")
title(s, "Project Objectives")

objs = [
    "Study the problem of identifying medicines from images",
    "Develop a Computer Vision and OCR pipeline",
    "Handle the Arabic and French text printed on the boxes",
    "Match the extracted information against a reference database",
    "Expose the final system through an API",
]
y0 = Inches(1.7)
for i, o in enumerate(objs):
    cy = y0 + Emu(int(Inches(1.02) * i))
    card(s, Inches(0.9), cy, Inches(11.5), Inches(0.82))
    numbered(s, Inches(1.15), cy + Inches(0.135), i + 1, diam=Inches(0.55))
    text(s, Inches(2.0), cy, Inches(10.2), Inches(0.82),
         [[(o, {"size": 16.5, "color": WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 4)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — STATE OF THE ART / METHODS
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "03 — State of the Art")
title(s, "Methods Studied", "Comparing the tools before making our choices")

# OCR column
card(s, Inches(0.9), Inches(1.85), Inches(5.6), Inches(4.9))
text(s, Inches(1.2), Inches(2.05), Inches(5.0), Inches(0.5),
     [[("OCR — reading the text", {"size": 18, "color": CYAN, "bold": True})]])
ocr_rows = [
    ("Tesseract", "Open-source reference, but unreliable for Arabic"),
    ("PaddleOCR", "Powerful but heavy to configure"),
    ("EasyOCR  ✓", "Native Arabic, simple, accepts images & paths"),
]
for i, (t, d) in enumerate(ocr_rows):
    cy = Inches(2.7) + Emu(int(Inches(1.3) * i))
    chosen = "✓" in t
    text(s, Inches(1.2), cy, Inches(5.0), Inches(1.2),
         [[(t, {"size": 15.5, "color": GREEN if chosen else WHITE, "bold": True})],
          [(d, {"size": 13, "color": MUTED})]], line_spacing=1.1)

# Recognition column
card(s, Inches(6.85), Inches(1.85), Inches(5.6), Inches(4.9))
text(s, Inches(7.15), Inches(2.05), Inches(5.0), Inches(0.5),
     [[("Identifying the name", {"size": 18, "color": CYAN, "bold": True})]])
rec_rows = [
    ("CNN alone", "Confuses visually similar boxes"),
    ("OCR + fuzzy matching  ✓", "Reads the printed name: reliable and direct"),
    ("Hybrid approach  ✓", "OCR as primary, CNN as fallback"),
]
for i, (t, d) in enumerate(rec_rows):
    cy = Inches(2.7) + Emu(int(Inches(1.3) * i))
    chosen = "✓" in t
    text(s, Inches(7.15), cy, Inches(5.0), Inches(1.2),
         [[(t, {"size": 15.5, "color": GREEN if chosen else WHITE, "bold": True})],
          [(d, {"size": 13, "color": MUTED})]], line_spacing=1.1)
page_num(s, 5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE / PIPELINE
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "04 — Architecture")
title(s, "Recognition Pipeline", "7 steps, from the image to the bilingual result")

steps = [
    ("Image", "Photo of the box"),
    ("CNN", "Visual prediction (MobileNetV2)"),
    ("Preprocessing", "Adaptive OpenCV"),
    ("OCR", "EasyOCR — FR + AR"),
    ("Name", "Fuzzy matching (RapidFuzz)"),
    ("Database", "MySQL lookup"),
    ("Result", "Dosage + form (FR / AR)"),
]
# two rows: 4 + 3
positions = [(0.9 + 3.0 * i, 1.95) for i in range(4)] + \
            [(2.4 + 3.0 * i, 4.35) for i in range(3)]
for i, (sx, sy) in enumerate(positions):
    L, T = Inches(sx), Inches(sy)
    card(s, L, T, Inches(2.7), Inches(1.95))
    numbered(s, L + Inches(0.25), T + Inches(0.25), i + 1, diam=Inches(0.5))
    text(s, L + Inches(0.2), T + Inches(0.85), Inches(2.3), Inches(0.5),
         [[(steps[i][0], {"size": 15, "color": CYAN, "bold": True})]])
    text(s, L + Inches(0.2), T + Inches(1.25), Inches(2.35), Inches(0.6),
         [[(steps[i][1], {"size": 11.5, "color": MUTED})]], line_spacing=1.05)
page_num(s, 6)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TECHNOLOGIES
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "05 — Technologies")
title(s, "Technical Stack")

techs = [
    ("Language", "Python 3.11"),
    ("CNN", "TensorFlow / Keras — MobileNetV2"),
    ("Image processing", "OpenCV"),
    ("OCR", "EasyOCR (PyTorch backend, GPU)"),
    ("Fuzzy matching", "RapidFuzz"),
    ("Database", "MySQL"),
    ("REST API", "FastAPI + Uvicorn"),
    ("Interface", "HTML / CSS / JavaScript"),
]
for i, (area, tech) in enumerate(techs):
    col = i % 2
    row = i // 2
    L = Inches(0.9 + col * 5.85)
    T = Inches(1.7 + row * 1.27)
    card(s, L, T, Inches(5.55), Inches(1.05))
    text(s, L + Inches(0.35), T + Inches(0.16), Inches(5.0), Inches(0.4),
         [[(area, {"size": 13, "color": CYAN, "bold": True})]])
    text(s, L + Inches(0.35), T + Inches(0.55), Inches(5.0), Inches(0.4),
         [[(tech, {"size": 15, "color": WHITE})]])
page_num(s, 7)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RECOGNITION STRATEGY (key decision)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "06 — Key Decision")
title(s, "Recognition: CNN + OCR")

card(s, Inches(0.9), Inches(1.75), Inches(5.6), Inches(2.1), CARD)
text(s, Inches(1.2), Inches(1.95), Inches(5.0), Inches(1.8),
     [[("The CNN problem", {"size": 17, "color": AMBER, "bold": True})],
      [("Trained on few images (~30/class), it confused Ventoline and "
        "Augmentin — it only sees the visual appearance.", {"size": 14, "color": WHITE})]],
     line_spacing=1.2)

card(s, Inches(6.85), Inches(1.75), Inches(5.6), Inches(2.1), CARD)
text(s, Inches(7.15), Inches(1.95), Inches(5.0), Inches(1.8),
     [[("Our solution", {"size": 17, "color": GREEN, "bold": True})],
      [("The name is PRINTED on the box → we read it with OCR + RapidFuzz. "
        "It is the most reliable source.", {"size": 14, "color": WHITE})]],
     line_spacing=1.2)

card(s, Inches(0.9), Inches(4.1), Inches(11.55), Inches(2.55), CARD2)
text(s, Inches(1.25), Inches(4.32), Inches(10.9), Inches(0.5),
     [[("Hybrid system — the right tool for each task",
        {"size": 17, "color": CYAN, "bold": True})]])
text(s, Inches(1.25), Inches(4.95), Inches(10.9), Inches(1.6),
     [[("•  OCR = ", {"size": 15, "color": WHITE, "bold": True}),
       ("primary source of the name (text read on the box)", {"size": 15, "color": MUTED})],
      [("•  CNN = ", {"size": 15, "color": WHITE, "bold": True}),
       ("fallback when the text is unreadable + rejects non-medicines", {"size": 15, "color": MUTED})],
      [("•  RapidFuzz = ", {"size": 15, "color": WHITE, "bold": True}),
       ("tolerates OCR typos (threshold 82%, or 65% with a 12-point margin)", {"size": 15, "color": MUTED})]],
     line_spacing=1.45)
page_num(s, 8)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MULTILINGUAL OCR & DOSAGE
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "07 — Multilingual OCR")
title(s, "Text Reading & Dosage Extraction")

blocks = [
    ("Two OCR engines", "EasyOCR does not mix Arabic and Latin: two separate "
     "readers (FR then AR+EN), then the results are merged."),
    ("Arabic units", "مجم، ملغ، ملجم → mg   •   ميكروغرام → mcg   •   مل → ml"),
    ("Arabic digits", "Two Unicode systems handled: ٠-٩ (Arabic) and ۰-۹ (Persian/Urdu) → 0-9"),
    ("Split number + unit", "Smart fallback when OCR separates « 500 » and « mg » on the box"),
]
for i, (h, d) in enumerate(blocks):
    cy = Inches(1.7) + Emu(int(Inches(1.27) * i))
    card(s, Inches(0.9), cy, Inches(11.55), Inches(1.07))
    text(s, Inches(1.25), cy + Inches(0.13), Inches(10.9), Inches(0.4),
         [[(h, {"size": 15.5, "color": CYAN, "bold": True})]])
    text(s, Inches(1.25), cy + Inches(0.55), Inches(10.9), Inches(0.45),
         [[(d, {"size": 13.5, "color": WHITE})]])
page_num(s, 9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ADAPTIVE PREPROCESSING
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "08 — OpenCV")
title(s, "Adaptive Preprocessing", "Every image is analyzed, then corrected only if needed")

# measure card
card(s, Inches(0.9), Inches(1.95), Inches(3.4), Inches(4.4), CARD2)
text(s, Inches(1.2), Inches(2.2), Inches(2.9), Inches(0.5),
     [[("1. Measure", {"size": 17, "color": CYAN, "bold": True})]])
for i, m in enumerate(["Brightness\n(mean)", "Contrast\n(std. deviation)",
                       "Sharpness\n(Laplacian variance)"]):
    text(s, Inches(1.2), Inches(2.85) + Emu(int(Inches(1.1) * i)), Inches(2.9), Inches(1.0),
         [[(m.split("\n")[0], {"size": 14, "color": WHITE, "bold": True})],
          [(m.split("\n")[1], {"size": 12, "color": MUTED})]], line_spacing=1.05)

# arrow
text(s, Inches(4.35), Inches(3.9), Inches(0.6), Inches(0.6),
     [[("→", {"size": 30, "color": CYAN, "bold": True})]], align=PP_ALIGN.CENTER)

# corrections
corr = [
    ("CLAHE", "Dark / low-contrast image → local contrast (LAB space)"),
    ("Unsharp mask", "Blurry image → edge enhancement"),
    ("Bilateral filter", "Noisy image → smoothing without blurring the text"),
]
text(s, Inches(5.1), Inches(2.0), Inches(7.0), Inches(0.5),
     [[("2. Correct (only if needed)", {"size": 17, "color": CYAN, "bold": True})]])
for i, (h, d) in enumerate(corr):
    cy = Inches(2.7) + Emu(int(Inches(1.25) * i))
    card(s, Inches(5.1), cy, Inches(7.35), Inches(1.05))
    text(s, Inches(5.4), cy + Inches(0.14), Inches(6.7), Inches(0.4),
         [[(h, {"size": 15, "color": GREEN, "bold": True})]])
    text(s, Inches(5.4), cy + Inches(0.55), Inches(6.8), Inches(0.45),
         [[(d, {"size": 13, "color": MUTED})]])
page_num(s, 10)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DATABASE
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "09 — Database")
title(s, "MySQL Reference Database")

stats = [("5", "medicines"), ("33", "variants"), ("2", "languages (FR / AR)")]
for i, (n, lbl) in enumerate(stats):
    L = Inches(0.9 + i * 3.95)
    card(s, L, Inches(1.8), Inches(3.6), Inches(1.7), CARD2)
    text(s, L, Inches(1.95), Inches(3.6), Inches(0.9),
         [[(n, {"size": 46, "color": CYAN, "bold": True})]], align=PP_ALIGN.CENTER)
    text(s, L, Inches(2.95), Inches(3.6), Inches(0.4),
         [[(lbl, {"size": 14, "color": MUTED})]], align=PP_ALIGN.CENTER)

card(s, Inches(0.9), Inches(3.8), Inches(11.55), Inches(2.85), CARD)
text(s, Inches(1.25), Inches(4.0), Inches(10.9), Inches(0.5),
     [[("Each variant stores:", {"size": 16, "color": CYAN, "bold": True})]])
cols = [
    ["name — name (French)", "name_ar — name (Arabic)", "dosage — e.g. 1000mg"],
    ["form — form (French)", "form_ar — form (Arabic)", "usage_fr / usage_ar — indication"],
]
for c, items in enumerate(cols):
    for r, it in enumerate(items):
        text(s, Inches(1.25 + c * 5.7), Inches(4.65) + Emu(int(Inches(0.62) * r)),
             Inches(5.5), Inches(0.5),
             [[("•  ", {"size": 14, "color": CYAN}),
               (it, {"size": 14, "color": WHITE})]])
page_num(s, 11)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — THREE-LEVEL VALIDATION
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "10 — Robustness")
title(s, "Validation: 3 Honest Verdicts", "The system never invents an identity")

verdicts = [
    (GREEN, "RECOGNIZED", "The name is read and matches the database — dosage and form identified"),
    (AMBER, "UNRECOGNIZED", "Text read clearly, but no known medicine matches"),
    (RED,   "REJECTED", "Not a medicine (e.g. animal, object) or unreadable text"),
]
for i, (col, h, d) in enumerate(verdicts):
    L = Inches(0.9 + i * 3.95)
    card(s, L, Inches(1.95), Inches(3.6), Inches(3.6), CARD)
    circle(s, L + Inches(1.4), Inches(2.3), Inches(0.8), col)
    text(s, L, Inches(3.35), Inches(3.6), Inches(0.5),
         [[(h, {"size": 19, "color": col, "bold": True})]], align=PP_ALIGN.CENTER)
    text(s, L + Inches(0.3), Inches(4.0), Inches(3.0), Inches(1.4),
         [[(d, {"size": 13.5, "color": MUTED})]], align=PP_ALIGN.CENTER, line_spacing=1.2)

text(s, Inches(0.9), Inches(5.85), Inches(11.55), Inches(0.9),
     [[("Safety: ", {"size": 15, "color": CYAN, "bold": True}),
       ("an “exact” variant is announced only if the dosage was actually "
        "read on the box — never inferred.", {"size": 15, "color": WHITE})]],
     line_spacing=1.2)
page_num(s, 12)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — API & INTERFACE
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "11 — Deliverable")
title(s, "REST API & Web Interface")

card(s, Inches(0.9), Inches(1.85), Inches(5.6), Inches(4.8), CARD)
text(s, Inches(1.2), Inches(2.05), Inches(5.0), Inches(0.5),
     [[("REST API — FastAPI", {"size": 18, "color": CYAN, "bold": True})]])
api = [
    ("POST /recognize", "Send an image → receive a complete JSON"),
    ("GET /", "Service health check"),
    ("Documentation /docs", "Interactive Swagger UI, auto-generated"),
    ("CORS enabled", "Usable by any external application"),
]
for i, (t, d) in enumerate(api):
    cy = Inches(2.7) + Emu(int(Inches(0.97) * i))
    text(s, Inches(1.2), cy, Inches(5.1), Inches(0.9),
         [[(t, {"size": 14.5, "color": GREEN, "bold": True, "font": "Consolas"})],
          [(d, {"size": 12.5, "color": MUTED})]], line_spacing=1.05)

card(s, Inches(6.85), Inches(1.85), Inches(5.6), Inches(4.8), CARD)
text(s, Inches(7.15), Inches(2.05), Inches(5.0), Inches(0.5),
     [[("Web interface", {"size": 18, "color": CYAN, "bold": True})]])
ui = [
    "Drag-and-drop a photo",
    "Clear result card (FR + AR)",
    "Colored status badge (recognized / rejected…)",
    "Confidence bar",
    "Collapsible technical details (OCR, OpenCV)",
]
for i, u in enumerate(ui):
    text(s, Inches(7.15), Inches(2.75) + Emu(int(Inches(0.72) * i)), Inches(5.1), Inches(0.6),
         [[("•  ", {"size": 15, "color": CYAN}), (u, {"size": 14.5, "color": WHITE})]])
page_num(s, 13)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DATASET
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "12 — Data")
title(s, "Dataset")

text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.0),
     [[("An image dataset was built to train the CNN, organized into folders "
        "by class (one folder = one class).", {"size": 16, "color": WHITE})]],
     line_spacing=1.25)

ds = [
    ("5 medicine classes", "Augmentin, Brufen, Doliprane, Smecta, Ventoline"),
    ("“other” class", "Non-medicine images, to learn to reject"),
    ("~30 images / class", "Small dataset → reason for the OCR-primary choice"),
    ("224×224 input", "Format expected by MobileNetV2"),
]
for i, (h, d) in enumerate(ds):
    col, row = i % 2, i // 2
    L = Inches(0.9 + col * 5.85)
    T = Inches(2.9 + row * 1.6)
    card(s, L, T, Inches(5.55), Inches(1.35))
    text(s, L + Inches(0.35), T + Inches(0.2), Inches(5.0), Inches(0.4),
         [[(h, {"size": 15.5, "color": CYAN, "bold": True})]])
    text(s, L + Inches(0.35), T + Inches(0.68), Inches(5.0), Inches(0.5),
         [[(d, {"size": 13, "color": MUTED})]])
page_num(s, 14)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — EVALUATION (headline numbers)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "13 — Experimental Evaluation")
title(s, "Results", "20 fresh images: medicines, blurry, Arabic, random objects")

big = [
    ("100%", "known medicines\ncorrectly named", GREEN, "11 / 11"),
    ("85%", "no false identity\never shown", CYAN, "17 / 20"),
    ("100%", "non-medicines\ncorrectly rejected", GREEN, "2 / 2"),
]
for i, (n, lbl, col, sub) in enumerate(big):
    L = Inches(0.9 + i * 3.95)
    card(s, L, Inches(2.0), Inches(3.6), Inches(2.9), CARD2)
    text(s, L, Inches(2.25), Inches(3.6), Inches(1.0),
         [[(n, {"size": 52, "color": col, "bold": True})]], align=PP_ALIGN.CENTER)
    text(s, L, Inches(3.35), Inches(3.6), Inches(0.8),
         [[(lbl.split("\n")[0], {"size": 14, "color": WHITE})],
          [(lbl.split("\n")[1], {"size": 14, "color": WHITE})]],
         align=PP_ALIGN.CENTER, line_spacing=1.05)
    text(s, L, Inches(4.35), Inches(3.6), Inches(0.4),
         [[(sub, {"size": 13, "color": MUTED, "italic": True})]], align=PP_ALIGN.CENTER)

text(s, Inches(0.9), Inches(5.3), Inches(11.55), Inches(1.4),
     [[("Highlights: ", {"size": 15, "color": CYAN, "bold": True}),
       ("Arabic boxes recognized at 100% • blurry Smecta → exact match • "
        "sideways Doliprane recognized at 94%.", {"size": 14.5, "color": WHITE})]],
     line_spacing=1.3)
page_num(s, 15)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — LIMITATIONS / ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "13 — Experimental Evaluation")
title(s, "Limitations & Failure Analysis")

lims = [
    ("Codoliprane → Doliprane", "Names 90% similar. Fix: add Codoliprane to "
     "the database. Illustrates the importance of a complete reference base."),
    ("Flipped box → rejected", "OCR cannot read upside-down text. "
     "Safe behavior: reject rather than guess."),
    ("“Medicine vs non-medicine”", "The candy box shows “99g” (net weight) — it "
     "looks more like a medicine than a poorly-photographed real box. "
     "Hard to tell apart from text alone."),
]
for i, (h, d) in enumerate(lims):
    cy = Inches(1.7) + Emu(int(Inches(1.6) * i))
    card(s, Inches(0.9), cy, Inches(11.55), Inches(1.4))
    circle(s, Inches(1.2), cy + Inches(0.45), Inches(0.5), AMBER)
    text(s, Inches(1.2), cy + Inches(0.45), Inches(0.5), Inches(0.5),
         [[("!", {"size": 22, "color": BG, "bold": True})]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.95), cy + Inches(0.18), Inches(10.2), Inches(0.4),
         [[(h, {"size": 15.5, "color": AMBER, "bold": True})]])
    text(s, Inches(1.95), cy + Inches(0.6), Inches(10.3), Inches(0.7),
         [[(d, {"size": 13, "color": MUTED})]], line_spacing=1.1)
page_num(s, 16)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — CONCLUSION & PERSPECTIVES
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "14 — Conclusion")
title(s, "Conclusion & Perspectives")

card(s, Inches(0.9), Inches(1.85), Inches(5.6), Inches(4.7), CARD)
text(s, Inches(1.2), Inches(2.05), Inches(5.0), Inches(0.5),
     [[("What works", {"size": 18, "color": GREEN, "bold": True})]])
done = [
    "Complete CV + OCR pipeline",
    "Arabic & French text handled",
    "Matching against a MySQL database",
    "REST API + web interface",
    "Robust to blur and rotation",
    "Honest 3-level validation",
]
for i, d in enumerate(done):
    text(s, Inches(1.2), Inches(2.7) + Emu(int(Inches(0.62) * i)), Inches(5.1), Inches(0.5),
         [[("✓  ", {"size": 14, "color": GREEN, "bold": True}),
           (d, {"size": 14, "color": WHITE})]])

card(s, Inches(6.85), Inches(1.85), Inches(5.6), Inches(4.7), CARD)
text(s, Inches(7.15), Inches(2.05), Inches(5.0), Inches(0.5),
     [[("Perspectives", {"size": 18, "color": CYAN, "bold": True})]])
todo = [
    "Expand the medicine database",
    "Grow the CNN dataset",
    "Orientation correction (rotation)",
    "Text-region detection",
    "Deploy the API online",
]
for i, d in enumerate(todo):
    text(s, Inches(7.15), Inches(2.7) + Emu(int(Inches(0.72) * i)), Inches(5.1), Inches(0.5),
         [[("→  ", {"size": 14, "color": CYAN, "bold": True}),
           (d, {"size": 14, "color": WHITE})]])
page_num(s, 17)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — THANKS
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
circle(s, Inches(-1.5), Inches(4.3), Inches(4.5), CARD)
circle(s, Inches(10.8), Inches(-1.6), Inches(4.5), CARD)
circle(s, Inches(11.7), Inches(-0.9), Inches(2.6), CYAN_D)

text(s, Inches(0), Inches(2.7), Inches(13.333), Inches(1.2),
     [[("Thank you for your attention", {"size": 44, "color": WHITE, "bold": True})]],
     align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(4.0), Inches(13.333), Inches(0.6),
     [[("Questions & discussion", {"size": 20, "color": CYAN})]], align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(5.5), Inches(13.333), Inches(0.5),
     [[("Hajar El Haj  &  Ikram Belmalhouz", {"size": 15, "color": MUTED})]],
     align=PP_ALIGN.CENTER)

prs.save("Presentation_Finale.pptx")
print("Saved Presentation_Finale.pptx —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
